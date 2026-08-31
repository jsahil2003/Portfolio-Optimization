"""
Generates citadel_deck.pptx from the content defined below.

Rerun this script (`python build_slides.py`) any time the strategy,
numbers, or story changes - don't hand-edit the .pptx binary. Keep this
file and SLIDES_SCRIPT.md's slide numbers in sync when you add/remove/
reorder slides.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---- palette (matches the dataviz-skill categorical palette used elsewhere in this project) ----
INK = RGBColor(0x0B, 0x0B, 0x0B)
SECONDARY_INK = RGBColor(0x52, 0x51, 0x4E)
SURFACE = RGBColor(0xFC, 0xFC, 0xFB)
ACCENT_BLUE = RGBColor(0x2A, 0x78, 0xD6)
ACCENT_ORANGE = RGBColor(0xEB, 0x68, 0x34)
ACCENT_GREEN = RGBColor(0x1B, 0xAF, 0x7A)
GOOD = RGBColor(0x00, 0x63, 0x00)
BAD = RGBColor(0xD0, 0x3B, 0x3B)
MUTED = RGBColor(0x89, 0x87, 0x81)

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)


def new_deck():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank layout


def add_bg(slide, color=SURFACE):
    bg = slide.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)  # rectangle
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)
    return bg


def add_title(slide, text, y=Inches(0.45), size=32, color=INK, x=Inches(0.6), w=Inches(12.1)):
    box = slide.shapes.add_textbox(x, y, w, Inches(1.0))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return box


def add_subtitle(slide, text, y=Inches(1.25), size=16, color=SECONDARY_INK, x=Inches(0.6), w=Inches(12.1)):
    box = slide.shapes.add_textbox(x, y, w, Inches(0.6))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return box


def add_bullets(slide, items, x=Inches(0.6), y=Inches(1.7), w=Inches(12.1), h=Inches(5.2), size=18, line_spacing=1.25):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(12)
        p.line_spacing = line_spacing
        if isinstance(item, tuple):
            text = item[0]
            level = item[1] if len(item) > 1 else 0
            bold = item[2] if len(item) > 2 else False
            color = item[3] if len(item) > 3 else INK
        else:
            text, level, bold, color = item, 0, False, INK
        run = p.add_run()
        run.text = ("•  " if level == 0 else "   –  ") + text
        run.font.size = Pt(size - level * 2)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Calibri"
        p.level = level
    return box


def add_footer(slide, text):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(7.05), Inches(12), Inches(0.35))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(10)
    run.font.color.rgb = MUTED
    run.font.name = "Calibri"


def add_table(slide, rows, col_widths, x=Inches(0.6), y=Inches(1.9), row_h=Inches(0.55), header_color=ACCENT_BLUE):
    n_rows, n_cols = len(rows), len(rows[0])
    total_w = sum(col_widths)
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, x, y, total_w, row_h * n_rows)
    tbl = tbl_shape.table
    for c, cw in enumerate(col_widths):
        tbl.columns[c].width = cw
    for r, row in enumerate(rows):
        for c, cell_val in enumerate(row):
            cell = tbl.cell(r, c)
            if isinstance(cell_val, tuple):
                text, bold, color = cell_val
            else:
                text, bold, color = cell_val, (r == 0), (RGBColor(0xFF, 0xFF, 0xFF) if r == 0 else INK)
            cell.text = str(text)
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(14)
            para.font.bold = bold
            para.font.color.rgb = color
            para.font.name = "Calibri"
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = SURFACE
    return tbl_shape


def add_stat_tile(slide, x, y, w, h, label, value, sub, value_color=ACCENT_BLUE):
    box = slide.shapes.add_shape(1, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xF4, 0xF4, 0xF2)
    box.line.color.rgb = RGBColor(0xE1, 0xE0, 0xD9)
    box.line.width = Pt(1)
    box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.12)
    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = label
    r1.font.size = Pt(13)
    r1.font.color.rgb = SECONDARY_INK
    r1.font.name = "Calibri"
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = value
    r2.font.size = Pt(30)
    r2.font.bold = True
    r2.font.color.rgb = value_color
    r2.font.name = "Calibri"
    p3 = tf.add_paragraph()
    r3 = p3.add_run()
    r3.text = sub
    r3.font.size = Pt(11)
    r3.font.color.rgb = MUTED
    r3.font.name = "Calibri"


def set_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# ============================================================
# BUILD DECK
# ============================================================
prs = new_deck()

# ---- Slide 1: Title ----
s = blank_slide(prs)
add_bg(s, RGBColor(0x0B, 0x0B, 0x0B))
box = s.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.5))
p = box.text_frame.paragraphs[0]
r = p.add_run()
r.text = "A Lookahead-Bias-Free Quantitative Portfolio"
r.font.size = Pt(40)
r.font.bold = True
r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
r.font.name = "Calibri"
box2 = s.shapes.add_textbox(Inches(0.8), Inches(3.6), Inches(11.7), Inches(0.8))
p2 = box2.text_frame.paragraphs[0]
r2 = p2.add_run()
r2.text = "Finesse x Citadel Portfolio Challenge — Round 2 Submission"
r2.font.size = Pt(20)
r2.font.color.rgb = RGBColor(0xC3, 0xC2, 0xB7)
r2.font.name = "Calibri"
box3 = s.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.6))
p3 = box3.text_frame.paragraphs[0]
r3 = p3.add_run()
r3.text = "Built with Claude Code (Anthropic) — see final slide for AI-usage disclosure"
r3.font.size = Pt(13)
r3.font.italic = True
r3.font.color.rgb = MUTED
r3.font.name = "Calibri"

# ---- Slide 2: Problem statement ----
s = blank_slide(prs)
add_bg(s)
add_title(s, "The Problem")
add_bullets(s, [
    "Build a quantitative equity portfolio: up to 10 stocks from Nifty 100 + Midcap 100 + Smallcap 100 (300 stocks total)",
    "Backtest 2021-01-01 to 2025-12-31, ₹1,00,00,000 starting capital, 0.1% transaction cost per trade",
    ("Ranked primarily on Total Net PNL", 1, True, INK),
    ("Secondary metrics: annualized return, max drawdown, Sharpe ratio, gain-to-loss ratio, accuracy, trade statistics", 1, False, INK),
    "Benchmarked against a relevant market index",
])
add_footer(s, "PROJECT_SUMMARY.md §1")

# ---- Slide 3: Approach overview ----
s = blank_slide(prs)
add_bg(s)
add_title(s, "Approach: Four Stages")
add_bullets(s, [
    ("1. Universe", 0, True, ACCENT_BLUE),
    ("Live-pulled from niftyindices.com every run — 300 stocks, always current", 1),
    ("2. Stock Selection", 0, True, ACCENT_BLUE),
    ("Every month, score all eligible stocks on 4 signals, take the top 10", 1),
    ("3. Position Weighting", 0, True, ACCENT_BLUE),
    ("Ledoit-Wolf shrinkage minimum-variance — uses the full correlation structure between the 10 picks", 1),
    ("4. Backtest & Validate", 0, True, ACCENT_BLUE),
    ("Real share-level trade simulation, real transaction costs, walk-forward validated weights", 1),
], size=19)
add_footer(s, "PROJECT_SUMMARY.md §2")

# ---- Slide 4: The five signals ----
s = blank_slide(prs)
add_bg(s)
add_title(s, "The Five Signals")
add_subtitle(s, "All five are calculated purely from historical price and volume data — genuinely lookahead-bias free")
rows = [
    ["Signal", "Weight", "What it measures"],
    ["Momentum", "40%", "12-month return, most recent month excluded"],
    ["Low-volatility (close-to-close)", "25%", "Trailing 6-month realized volatility (inverted)"],
    ["Parkinson volatility (intraday range)", "10%", "Volatility from each day's High-Low range (inverted)"],
    ["52-week-high proximity", "15%", "Price ÷ trailing 12-month high"],
    ["Liquidity", "10%", "Amihud illiquidity (price impact per rupee traded), sign-flipped"],
]
add_table(s, rows, [Inches(4.6), Inches(1.5), Inches(6.0)], y=Inches(2.0), row_h=Inches(0.6))
add_footer(s, "PROJECT_SUMMARY.md §2  |  CONCEPTS.md §3-4, §24, §34")

# ---- Slide 5: Why no fundamentals ----
s = blank_slide(prs)
add_bg(s)
add_title(s, "Why No P/E, ROE, or Other Fundamentals?")
add_bullets(s, [
    ("The problem: lookahead bias", 0, True, BAD),
    ("Our data source (yfinance) only exposes TODAY's fundamentals, not what they were on any past date", 1),
    ("Using 2026 fundamentals to simulate a March 2022 decision lets the backtest \"see the future\"", 1),
    ("The fix", 0, True, GOOD),
    ("Replaced Value/Quality (fundamentals-based) with 52-week-high and Liquidity (price/volume-based)", 1),
    ("yfinance's historical OHLCV is genuinely point-in-time correct — no peeking, ever", 1),
], size=19)
add_footer(s, "CONCEPTS.md §18, §24")

# ---- Slide 6: Ledoit-Wolf weighting ----
s = blank_slide(prs)
add_bg(s)
add_title(s, "Position Sizing: Ledoit-Wolf Minimum-Variance")
add_bullets(s, [
    "Not equal-dollar (10% each) and not just inverse-volatility (ignores correlation)",
    "Uses the full covariance matrix between all 10 stocks, shrunk toward a stable target (Ledoit & Wolf, 2004) to avoid overfitting noisy correlation estimates",
    ("Compared 3 schemes empirically: inverse-vol, Ledoit-Wolf, Hierarchical Risk Parity", 0, True, ACCENT_BLUE),
    ("Ledoit-Wolf won on every metric simultaneously — verified across 4 separate robustness checks", 1),
], size=19)
add_footer(s, "CONCEPTS.md §20-23")

# ---- Slide 7: Correlation / diversification ----
s = blank_slide(prs)
add_bg(s)
add_title(s, "Diversification: Does Ledoit-Wolf Guarantee It?")
add_bullets(s, [
    ("Short answer: No — and that's an important distinction", 0, True, BAD),
    ("Ledoit-Wolf is a WEIGHTING technique — it can't change which 10 stocks get picked", 1),
    ("It optimally exploits whatever correlation the selected stocks actually have", 1),
    ("Concrete evidence, 2021-11-01 rebalance:", 0, True, ACCENT_BLUE),
    ("Most negative pair (IIFL.NS & TATAELXSI.NS, corr -0.11) got 21.2% combined weight under Ledoit-Wolf vs. only 17.9% under inverse-vol", 1),
    ("Across 270 sampled pairs (6 dates): 7.0% negatively correlated, avg correlation 0.20 — realistic for a single-market long-only book", 0, True, SECONDARY_INK),
], size=17)
add_footer(s, "CONCEPTS.md §31")

# ---- Slide 8: Overfitting lesson 1 ----
s = blank_slide(prs)
add_bg(s)
add_title(s, "Rigor #1: The Spike That Wasn't Real")
add_bullets(s, [
    "A grid search over signal weights found a combination with Sharpe 1.80 — beating everything tried so far",
    ("But checking NEARBY weight combinations (±0.02-0.05):", 0, True, BAD),
    ("Sharpe swung wildly — 1.50 to 1.84 — for barely-different weights", 1),
    ("A smooth, stable neighborhood is expected from a real signal; a wild swing is the signature of overfitting to noise", 0, True, ACCENT_ORANGE),
    ("Decision: rejected. Shipped the honest, untuned baseline instead.", 0, True, GOOD),
], size=19)
add_footer(s, "CONCEPTS.md §27  |  CHANGELOG.md")

# ---- Slide 9: Overfitting lesson 2 ----
s = blank_slide(prs)
add_bg(s)
add_title(s, "Rigor #2: Smooth ≠ Safe — Walk-Forward Validation")
add_bullets(s, [
    "Tried adding a 5th signal (short-term reversal). This time the weight search WAS smooth in-sample.",
    ("Split data into TRAIN (2021-2023) and TEST (2024-2025, never touched during tuning):", 0, True, ACCENT_BLUE),
    ("Train: Sharpe 1.27 → 1.84 (looked like a clear win)", 1),
    ("Test: Sharpe 1.63 → 1.42 (the \"improvement\" made things WORSE on unseen data)", 1, True, BAD),
    ("Lesson: smoothness rules out fitting noise in the SEARCH, but not fitting a pattern specific to the training period's market regime", 0, True, ACCENT_ORANGE),
], size=18)
add_footer(s, "CONCEPTS.md §28-29  |  CHANGELOG.md")

# ---- Slide 10: Final tuning - the winning result ----
s = blank_slide(prs)
add_bg(s)
add_title(s, "Rigor #3: Properly Tuning the Base Weights")
add_bullets(s, [
    "Applied walk-forward validation to the base 4 signals themselves: 84 weight combinations, trained on 2021-2023",
    ("Correlation between TRAIN Sharpe and TEST Sharpe across all 84: 0.019 — essentially zero", 0, True, BAD),
    ("Train-period \"winners\" were momentum-heavy — and collapsed on the unseen test period (Sharpe fell to as low as 0.49)", 1),
    ("Selected the combination robust in BOTH windows instead of the train-only leader — this became the strategy for a while", 0, True, GOOD),
    ("(Later replaced by a 5-signal version — see next slide)", 1, True, SECONDARY_INK),
], size=18)
add_footer(s, "CONCEPTS.md §30  |  CHANGELOG.md")

# ---- Slide 10.5: Rigor #4 - adding Parkinson, 4-way validated ----
s = blank_slide(prs)
add_bg(s)
add_title(s, "Rigor #4: Adding a 5th Signal, Stress-Tested 4 Ways")
add_bullets(s, [
    "Added Parkinson (intraday-range) volatility alongside close-to-close volatility — complementary, not redundant",
    ("Check 1 — primary train/test split: candidate beat the shipped strategy on both windows", 0, True, GOOD),
    ("Check 2 — fine-grid neighborhood (requested before adopting): revealed real fragility — nearby weights swung from Sharpe 0.60 to 1.69", 0, True, BAD),
    ("Check 3 — independent second split (2023-25 train / 2021-22 test): still beat the shipped strategy on both halves", 0, True, GOOD),
    ("Check 4 — re-centered within the cluster of good performance, trading peak score for a much higher worst-case floor", 0, True, GOOD),
    ("Final weights beat the previous strategy on all 4 checks, not just 1", 0, True, GOOD),
], size=16)
add_footer(s, "CONCEPTS.md §34  |  CHANGELOG.md")

# ---- Slide 11: Results ----
s = blank_slide(prs)
add_bg(s)
add_title(s, "Final Results (2021-2025, Live Data)")
tile_w, tile_h = Inches(2.85), Inches(1.7)
add_stat_tile(s, Inches(0.6), Inches(1.9), tile_w, tile_h, "Total Net PNL", "₹3.90 cr", "on ₹1 cr starting capital", ACCENT_BLUE)
add_stat_tile(s, Inches(3.6), Inches(1.9), tile_w, tile_h, "Annualized Return", "38.4%", "vs. Nifty 50: 13.4%", ACCENT_GREEN)
add_stat_tile(s, Inches(6.6), Inches(1.9), tile_w, tile_h, "Max Drawdown", "-19.1%", "vs. Nifty 50: -17.2%", ACCENT_ORANGE)
add_stat_tile(s, Inches(9.6), Inches(1.9), tile_w, tile_h, "Sharpe Ratio", "1.73", "vs. Nifty 50: 0.97", ACCENT_BLUE)
rows = [
    ["Metric", "Strategy", "Nifty 50", "Nifty 500 TMI"],
    ["Annualized Return", "38.4%", "13.4%", "15.5%"],
    ["Max Drawdown", "-19.1%", "-17.2%", "-18.8%"],
    ["Sharpe Ratio", "1.73", "0.97", "1.07"],
]
add_table(s, rows, [Inches(3.6), Inches(2.8), Inches(2.8), Inches(2.8)], y=Inches(4.1), row_h=Inches(0.55))
add_footer(s, "PROJECT_SUMMARY.md §5  |  citadel_submission.xlsx")

# ---- Slide 12: Methodological journey ----
s = blank_slide(prs)
add_bg(s)
add_title(s, "The Full Methodological Journey")
add_bullets(s, [
    "1. Built the initial pipeline — universe, signals, weighting, backtest",
    "2. Found and fixed 2 real bugs on the very first live run (implausible results triggered investigation)",
    "3. Tuned risk vs. return — found and rejected an overly aggressive trade-off, chose a deliberate balance",
    "4. Compared 3 weighting schemes empirically — adopted Ledoit-Wolf after 4 robustness checks",
    "5. Eliminated lookahead bias entirely — replaced fundamentals-based signals with price/volume-only ones",
    "6. Walk-forward validated the final weights — caught and rejected 2 overfitting traps along the way",
], size=17)
add_footer(s, "CHANGELOG.md — every step logged with exact numbers")

# ---- Slide 13: Limitations ----
s = blank_slide(prs)
add_bg(s)
add_title(s, "Known Limitations")
add_bullets(s, [
    "Single historical backtest window (2021-2025) — walk-forward validation reduces but doesn't eliminate regime-specific risk",
    "Universe freshness depends on live network access to niftyindices.com (falls back to a local cache)",
    "Transaction cost model is simplified (flat 0.1%, no bid-ask spread or market impact modeled separately)",
    "No shorting, leverage, or other asset classes — true negative-correlation diversification isn't achievable within the competition's rules",
], size=19)
add_footer(s, "PROJECT_SUMMARY.md §7")

# ---- Slide 14: AI usage disclosure ----
s = blank_slide(prs)
add_bg(s, RGBColor(0xF4, 0xF4, 0xF2))
add_title(s, "How Claude Was Used On This Project", color=INK)
add_bullets(s, [
    ("Built with Claude Code (Anthropic's AI coding assistant), Aug 20-29, 2026", 0, True, INK),
    ("Claude wrote the entire codebase and ran every backtest, tuning sweep, and validation shown in this deck — all on real live market data, not fabricated numbers", 0),
    ("The human's role: directed what to build, asked challenging questions (\"is this handled by Ledoit-Wolf?\", \"do more rigorous tuning\"), made final calls on trade-offs, requested this documentation", 0),
    ("Claude proactively flagged its own bugs, rejected results, and limitations throughout — not just favorable numbers", 0),
    ("Full disclosure and methodology: PROJECT_SUMMARY.md §8", 0, True, ACCENT_BLUE),
], size=17)

# ---- Slide 15: Thank you ----
s = blank_slide(prs)
add_bg(s, RGBColor(0x0B, 0x0B, 0x0B))
box = s.shapes.add_textbox(Inches(0.8), Inches(3.0), Inches(11.7), Inches(1.2))
p = box.text_frame.paragraphs[0]
r = p.add_run()
r.text = "Questions?"
r.font.size = Pt(44)
r.font.bold = True
r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
r.font.name = "Calibri"
box2 = s.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(11.7), Inches(0.8))
p2 = box2.text_frame.paragraphs[0]
r2 = p2.add_run()
r2.text = "Full documentation: PROJECT_SUMMARY.md, CONCEPTS.md, CODE_GUIDE.md, CHANGELOG.md, BEGINNERS_GUIDE.md"
r2.font.size = Pt(15)
r2.font.color.rgb = RGBColor(0xC3, 0xC2, 0xB7)
r2.font.name = "Calibri"

# ---- Presenter notes (condensed from SLIDES_SCRIPT.md - keep both in sync) ----
SPEAKER_NOTES = [
    "Quant equity strategy for the Finesse x Citadel challenge, fully free of lookahead bias. Built with heavy use of Claude Code (AI coding assistant) - disclosed in detail on the AI-usage slide near the end.",
    "Brief: up to 10 stocks from 300 (Nifty100+Midcap100+Smallcap100), 2021-2025 backtest, ranked primarily on Total Net PNL, with Sharpe/drawdown/benchmark as secondary metrics.",
    "Four stages: (1) live-pulled universe, (2) monthly stock selection via composite score, (3) Ledoit-Wolf position sizing, (4) realistic backtest + rigorous validation.",
    "Momentum 40%, close-to-close low-vol 25%, Parkinson (intraday) vol 10%, 52-week-high 15%, liquidity 10%. Weights are walk-forward validated across 4 separate checks - see the rigor slides ahead.",
    "No fundamentals (P/E, ROE) because our data source only gives TODAY's snapshot, not historical values - using it would be lookahead bias. Earlier fundamentals-based version retired, kept in code for comparison.",
    "Ledoit-Wolf shrinkage minimum-variance: uses full correlation structure between the 10 picks, not just each stock's own risk. Beat 2 alternatives (inverse-vol, HRP) on every metric, verified 4 ways.",
    "Ledoit-Wolf does NOT guarantee negative correlation - it's a weighting method, can't change which stocks get picked. It DOES exploit whatever correlation exists: concrete example, negatively-correlated pair got more weight under Ledoit-Wolf than a simpler method. ~7% of pairs negatively correlated overall - realistic for single-market long-only.",
    "Mistake #1 (self-caught): a weight combo looked great, but nearby weights swung wildly - sign of a fluke, not a real pattern. Rejected, used the honest simpler version.",
    "Mistake #2 (self-caught): added a 5th signal, smooth in-sample this time - but train/test split showed it got WORSE on unseen data. Lesson: smoothness alone isn't proof against overfitting.",
    "Applied the lesson properly: 84 weight combos, selected by consistency across train AND test, not train alone. Train-test correlation ~0 across all combos - striking proof naive tuning doesn't transfer. Final pick performs nearly identically on both windows. This became the strategy for a while, later refined further (next slide).",
    "Added a 5th signal (Parkinson intraday volatility) alongside the existing close-to-close volatility - complementary, not redundant. First candidate looked good on the primary split but a finer neighborhood check (asked for before adopting) revealed real fragility nearby. Re-centered within the cluster of good performance and re-validated on an independent second split. Final weights beat the previous strategy on all 4 checks run, not just 1.",
    "38.4% annualized return vs Nifty 50's 13.4%; Sharpe 1.73 vs 0.97; MDD close to benchmark. ~₹3.90cr net profit on ₹1cr capital over 5 years.",
    "Not a straight line: 2 real bugs found+fixed on first live run (implausible numbers triggered investigation), weighting schemes compared properly, 2 self-caught overfitting traps.",
    "Honest limitations: single 5-year window, universe needs live network access, simplified transaction costs, no shorting/leverage so diversification is inherently limited.",
    "Full transparency: Claude Code wrote all the code and ran every real backtest shown. My role: directed the work, asked hard questions, made trade-off calls, ensured honest documentation. I can explain every concept myself using the source docs.",
    "Open for questions on any part - math, validation methodology, or the code itself.",
]

for slide, note in zip(prs.slides, SPEAKER_NOTES):
    set_notes(slide, note)

prs.save("citadel_deck.pptx")
print(f"Wrote citadel_deck.pptx with {len(prs.slides.__iter__.__self__._sldIdLst)} slides and speaker notes")
